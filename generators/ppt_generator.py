import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
from utils.fonts import get_pil_font
from utils.colors import hex_to_rgb, get_text_color_for_background

def add_title_slide(prs, business_name, tagline, owner_name, brand_rgb):
    """Slide 1: Dark branding theme cover page."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background (Dark Slate)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # #0F172A Slate 900
    
    # Large brand color accent shape (top triangle/polygon or header strip)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*brand_rgb)
    shape.line.fill.background()
    
    # Text Box for Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Business name
    p = tf.paragraphs[0]
    p.text = business_name
    p.font.name = "Montserrat"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    
    # Tagline
    if tagline:
        p2 = tf.add_paragraph()
        p2.text = tagline
        p2.font.name = "Inter"
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(148, 163, 184)  # slate 400
        p2.space_before = Pt(20)
        p2.alignment = PP_ALIGN.LEFT
        
    # Owner & presenting details
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(10.5), Inches(1.2))
    tf2 = txBox2.text_frame
    
    p3 = tf2.paragraphs[0]
    p3.text = f"Presented by: {owner_name}"
    p3.font.name = "Inter"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(*brand_rgb)
    
    p4 = tf2.add_paragraph()
    p4.text = "Powered by Magic Business"
    p4.font.name = "Inter"
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(100, 116, 139)  # slate 500
    p4.space_before = Pt(5)

def add_content_slide(prs, title, brand_rgb):
    """Create a standardized content slide layout with brand colors."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)  # #F8FAFC slate 50
    
    # Top banner stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = RGBColor(*brand_rgb)
    stripe.line.fill.background()
    
    # Title Text Box
    titleBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.0))
    tf = titleBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Montserrat"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    
    # Divider line under title
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(2.5), Inches(0.04))
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(*brand_rgb)
    divider.line.fill.background()
    
    return slide

def add_about_slide(prs, description, brand_rgb):
    """Slide 2: About us."""
    slide = add_content_slide(prs, "About Us", brand_rgb)
    
    # Left border highlight shape
    highlight = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.08), Inches(4.0))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(*brand_rgb)
    highlight.line.fill.background()
    
    # Text Box
    textBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(4.0))
    tf = textBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = description
    p.font.name = "Inter"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(51, 65, 85)  # slate 700
    p.line_spacing = 1.3

def add_services_slide(prs, services, brand_rgb):
    """Slide 3: List of services."""
    slide = add_content_slide(prs, "Our Core Services", brand_rgb)
    
    # List services in two columns or grid depending on number
    left_col = Inches(0.8)
    right_col = Inches(6.8)
    top_y = Inches(2.2)
    width_col = Inches(5.5)
    
    for i, s in enumerate(services[:6]):
        col = left_col if i % 2 == 0 else right_col
        row = i // 2
        y = top_y + Inches(row * 1.3)
        
        # Draw a small brand color accent dot/bullet
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, col, y + Inches(0.08), Inches(0.18), Inches(0.18))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = RGBColor(*brand_rgb)
        bullet.line.fill.background()
        
        # Text box next to bullet
        box = slide.shapes.add_textbox(col + Inches(0.35), y, width_col, Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s.strip()
        p.font.name = "Montserrat"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)

def add_usps_slide(prs, brand_rgb):
    """Slide 4: Why Choose Us (USPs)."""
    slide = add_content_slide(prs, "Why Choose Us", brand_rgb)
    
    usps = [
        ("Premium Quality", "We maintain the highest standards of quality in all our deliverables."),
        ("Customer Focused", "Our services are tailored entirely to maximize client value and satisfaction."),
        ("Professional Team", "Work with experienced industry veterans who understand your requirements."),
        ("Affordable Rates", "Professional value metrics offered at competitive business friendly rates.")
    ]
    
    col_width = Inches(5.5)
    col1 = Inches(0.8)
    col2 = Inches(6.8)
    
    for i, (title, desc) in enumerate(usps):
        col = col1 if i % 2 == 0 else col2
        row = i // 2
        y = Inches(2.2) + Inches(row * 2.0)
        
        # Accent rounded border box for each USP card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col, y, col_width, Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Line border
        card.line.color.rgb = RGBColor(*brand_rgb)
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Montserrat"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*brand_rgb)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Inter"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(71, 85, 105)  # slate 600
        p2.space_before = Pt(8)

def add_pricing_slide(prs, services, brand_rgb):
    """Slide 5: Pricing table or call-to-action."""
    slide = add_content_slide(prs, "Service Pricing Guide", brand_rgb)
    
    # Parse prices if they exist in services, e.g. "Grocery - Rs. 100"
    items_with_price = []
    for s in services:
        if '-' in s or '|' in s or ':' in s:
            sep = '-' if '-' in s else ('|' if '|' in s else ':')
            parts = s.split(sep, 1)
            items_with_price.append((parts[0].strip(), parts[1].strip()))
            
    if len(items_with_price) >= 2:
        # Create a table
        rows = min(5, len(items_with_price) + 1)
        left = Inches(0.8)
        top = Inches(2.2)
        width = Inches(11.5)
        height = Inches(0.8 * rows)
        
        table_shape = slide.shapes.add_table(rows, 2, left, top, width, height)
        table = table_shape.table
        
        # Style headers
        table.columns[0].width = Inches(8.0)
        table.columns[1].width = Inches(3.5)
        
        cell_head_name = table.cell(0, 0)
        cell_head_name.text = "Service Package / Option"
        cell_head_name.fill.solid()
        cell_head_name.fill.fore_color.rgb = RGBColor(*brand_rgb)
        
        cell_head_price = table.cell(0, 1)
        cell_head_price.text = "Price"
        cell_head_price.fill.solid()
        cell_head_price.fill.fore_color.rgb = RGBColor(*brand_rgb)
        
        for idx, (name, val) in enumerate(items_with_price[:rows-1]):
            cell_n = table.cell(idx+1, 0)
            cell_n.text = name
            cell_p = table.cell(idx+1, 1)
            cell_p.text = val
            
            # Simple alternating background
            color = RGBColor(241, 245, 249) if idx % 2 == 0 else RGBColor(255, 255, 255)
            cell_n.fill.solid()
            cell_n.fill.fore_color.rgb = color
            cell_p.fill.solid()
            cell_p.fill.fore_color.rgb = color
    else:
        # Standard Call To Action (CTA) Pricing Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(2.2), Inches(8.33), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(*brand_rgb)
        card.line.width = Pt(2)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.4)
        
        p = tf.paragraphs[0]
        p.text = "Tailored Custom Quotations"
        p.font.name = "Montserrat"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "We offer customizable pricing packages to suit your business requirements perfectly. Reach out directly for a personalized quote, or explore custom contracts."
        p2.font.name = "Inter"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(100, 116, 139)
        p2.space_before = Pt(20)
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = tf.add_paragraph()
        p3.text = "Contact Us Today to Request a Quote"
        p3.font.name = "Montserrat"
        p3.font.size = Pt(20)
        p3.font.bold = True
        p3.font.color.rgb = RGBColor(*brand_rgb)
        p3.space_before = Pt(35)
        p3.alignment = PP_ALIGN.CENTER

def add_portfolio_slide(prs, brand_rgb):
    """Slide 6: Visual work blocks."""
    slide = add_content_slide(prs, "Our Work / Portfolio", brand_rgb)
    
    # 3 Placeholder card boxes
    width_box = Inches(3.6)
    height_box = Inches(3.8)
    top_y = Inches(2.2)
    
    for i in range(3):
        x = Inches(0.8) + Inches(i * 4.1)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top_y, width_box, height_box)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(241, 245, 249) # light grey
        card.line.color.rgb = RGBColor(203, 213, 225) # border
        
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Project Showcase {i+1}\n\n[Double click to insert project image here]"
        p.font.name = "Inter"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(148, 163, 184)
        p.alignment = PP_ALIGN.CENTER

def add_contact_slide(prs, owner_name, phone, email, address, website, brand_rgb):
    """Slide 7: Reach out page."""
    slide = add_content_slide(prs, "Get In Touch", brand_rgb)
    
    # Large details box
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.0), Inches(4.0))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    details = [
        ("Owner:", owner_name),
        ("Phone:", phone),
        ("Email:", email)
    ]
    if website:
        details.append(("Website:", website))
    if address:
        details.append(("Address:", address.replace('\n', ', ')))
        
    for i, (k, v) in enumerate(details):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"{k} "
        p.font.name = "Montserrat"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*brand_rgb)
        
        # Add value in normal text
        run = p.add_run()
        run.text = v
        run.font.name = "Inter"
        run.font.size = Pt(20)
        run.font.bold = False
        run.font.color.rgb = RGBColor(51, 65, 85)
        p.space_after = Pt(20)
        
    # Right box (Mock Map area)
    map_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(2.2), Inches(4.8), Inches(3.8))
    map_box.fill.solid()
    map_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    map_box.line.color.rgb = RGBColor(*brand_rgb)
    map_box.line.width = Pt(1.5)
    
    tf_m = map_box.text_frame
    p_m = tf_m.paragraphs[0]
    p_m.text = "[ Location Reference Grid ]\n\nShop Mockup Map Placeholder"
    p_m.font.name = "Inter"
    p_m.font.size = Pt(16)
    p_m.font.color.rgb = RGBColor(148, 163, 184)
    p_m.alignment = PP_ALIGN.CENTER

def add_thank_you_slide(prs, business_name, tagline, brand_rgb):
    """Slide 8: Dark slide thanking audience."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Dark Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    # Text Frame
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.name = "Montserrat"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*brand_rgb)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"We look forward to working with you."
    p2.font.name = "Inter"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(20)
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = f"{business_name}  |  {tagline}" if tagline else business_name
    p3.font.name = "Montserrat"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(*brand_rgb)
    p3.space_before = Pt(30)
    p3.alignment = PP_ALIGN.CENTER

def generate_ppt(
    business_name: str,
    owner_name: str,
    phone: str,
    email: str,
    website: str = "",
    address: str = "",
    tagline: str = "",
    description: str = "",
    services: list[str] = None,
    brand_color_hex: str = "#2563EB",
    output_dir: str = ""
) -> dict:
    """
    Generate professional pitch deck PPTX and preview thumbnail.
    """
    os.makedirs(output_dir, exist_ok=True)
    generated_files = {}
    
    brand_rgb = hex_to_rgb(brand_color_hex)
    if not services:
        services = ["General Consultations", "Premium Delivery Support", "Quality Client Services"]
        
    if not description:
        description = f"Welcome to {business_name}, spearheaded by {owner_name}. We deliver exceptional services customized to fit your dynamic workflow requirements."
        
    # Create Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Cover
    add_title_slide(prs, business_name, tagline, owner_name, brand_rgb)
    # Slide 2: About
    add_about_slide(prs, description, brand_rgb)
    # Slide 3: Services
    add_services_slide(prs, services, brand_rgb)
    # Slide 4: USPs
    add_usps_slide(prs, brand_rgb)
    # Slide 5: Pricing
    add_pricing_slide(prs, services, brand_rgb)
    # Slide 6: Portfolio
    add_portfolio_slide(prs, brand_rgb)
    # Slide 7: Contact
    add_contact_slide(prs, owner_name, phone, email, address, website, brand_rgb)
    # Slide 8: Thank You
    add_thank_you_slide(prs, business_name, tagline, brand_rgb)
    
    # Save PPTX
    pptx_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(pptx_path)
    generated_files["presentation_pptx"] = "presentation.pptx"
    
    # ---------------- GENERATE PILLOW PREVIEW PNG (Slide 1 Mockup) ----------------
    # Widescreen aspect ratio (16:9) -> e.g. 960 x 540 px
    pw, ph = 960, 540
    preview = Image.new("RGBA", (pw, ph), (15, 23, 42, 255))
    draw_p = ImageDraw.Draw(preview)
    
    # Brand color left highlight strip
    draw_p.rectangle([0, 0, 30, ph], fill=brand_rgb)
    
    # Title Text
    font_title = get_pil_font('montserrat_bold', 48)
    draw_p.text((100, 140), business_name, fill=(255, 255, 255), font=font_title)
    
    # Tagline Text
    if tagline:
        font_tag = get_pil_font('inter_regular', 20)
        draw_p.text((100, 220), tagline, fill=(148, 163, 184), font=font_tag)
        
    # Presenter Details
    font_owner = get_pil_font('inter_bold', 18)
    draw_p.text((100, ph - 140), f"Presented by: {owner_name}", fill=brand_rgb, font=font_owner)
    
    font_footer = get_pil_font('inter_regular', 12)
    draw_p.text((100, ph - 95), "Powered by Magic Business (Pitch Deck Preview)", fill=(100, 116, 139), font=font_footer)
    
    preview_path = os.path.join(output_dir, "presentation_preview.png")
    preview.convert("RGB").save(preview_path, "PNG")
    generated_files["presentation_preview"] = "presentation_preview.png"
    
    return generated_files
